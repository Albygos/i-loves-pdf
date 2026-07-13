import io
from pypdf import PdfReader, PdfWriter
import fitz  # PyMuPDF
from PIL import Image
from docx import Document
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter, landscape
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak, Table as RLTable, TableStyle as RLTableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors
from pptx import Presentation
from pptx.util import Inches, Pt
import openpyxl
from deep_translator import GoogleTranslator
from html.parser import HTMLParser

# --- Word Helpers ---
from docx.document import Document as Doc
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.table import Table as DocxTable
from docx.text.paragraph import Paragraph as DocxParagraph

def iter_block_items(parent):
    if isinstance(parent, Doc):
        parent_elm = parent.element.body
    else:
        parent_elm = parent._element

    for child in parent_elm.iterchildren():
        if isinstance(child, CT_P):
            yield DocxParagraph(child, parent)
        elif isinstance(child, CT_Tbl):
            yield DocxTable(child, parent)

# --- HTML Helpers ---
class ReportLabHTMLParser(HTMLParser):
    def __init__(self, styles, normal_style):
        super().__init__()
        self.styles = styles
        self.normal_style = normal_style
        self.story = []
        self.current_style = normal_style
        self.current_text = ""
        
    def handle_starttag(self, tag, attrs):
        if tag in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6']:
            self.flush()
            level = tag[1]
            if level == '1':
                self.current_style = self.styles['Heading1']
            elif level == '2':
                self.current_style = self.styles['Heading2']
            else:
                self.current_style = self.styles['Heading3']
        elif tag == 'p':
            self.flush()
            self.current_style = self.normal_style
        elif tag == 'br':
            self.current_text += "<br/>"
        elif tag in ['b', 'strong']:
            self.current_text += "<b>"
        elif tag in ['i', 'em']:
            self.current_text += "<i>"
        elif tag in ['u']:
            self.current_text += "<u>"
        elif tag == 'a':
            href = dict(attrs).get('href', '')
            self.current_text += f'<a href="{href}">'
            
    def handle_endtag(self, tag):
        if tag in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'p']:
            self.flush()
        elif tag in ['b', 'strong']:
            self.current_text += "</b>"
        elif tag in ['i', 'em']:
            self.current_text += "</i>"
        elif tag in ['u']:
            self.current_text += "</u>"
        elif tag == 'a':
            self.current_text += "</a>"
            
    def handle_data(self, data):
        safe_data = data.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
        self.current_text += safe_data
        
    def flush(self):
        text = self.current_text.strip()
        if text:
            try:
                self.story.append(Paragraph(text, self.current_style))
                self.story.append(Spacer(1, 8))
            except Exception:
                import re
                plain_text = re.sub('<[^<]+?>', '', text)
                self.story.append(Paragraph(plain_text, self.current_style))
                self.story.append(Spacer(1, 8))
        self.current_text = ""
        self.current_style = self.normal_style


# --- Converter Helper Methods ---

def convert_pptx_to_pdf(pptx_path, pdf_path):
    from reportlab.lib.units import inch
    prs = Presentation(pptx_path)
    slide_width = prs.slide_width / 12700 if prs.slide_width else 720
    slide_height = prs.slide_height / 12700 if prs.slide_height else 540
    
    c = canvas.Canvas(pdf_path, pagesize=(slide_width, slide_height))
    
    for slide in prs.slides:
        for shape in slide.shapes:
            left = shape.left / 12700 if shape.left else 0
            top = shape.top / 12700 if shape.top else 0
            width = shape.width / 12700 if shape.width else 100
            height = shape.height / 12700 if shape.height else 100
            
            if shape.has_text_frame:
                text = shape.text.strip()
                if text:
                    text_obj = c.beginText()
                    y_start = slide_height - top - 12
                    text_obj.setTextOrigin(left + 10, y_start)
                    text_obj.setFont("Helvetica", 12)
                    text_obj.setFillColorRGB(0.1, 0.1, 0.1)
                    
                    for paragraph in shape.text_frame.paragraphs:
                        p_text = paragraph.text.strip()
                        if p_text:
                            font_size = 12
                            if paragraph.font and paragraph.font.size:
                                font_size = paragraph.font.size / 12700
                            
                            text_obj.setFont("Helvetica-Bold" if (paragraph.font and paragraph.font.bold) else "Helvetica", font_size or 12)
                            text_obj.textLine(p_text)
                    c.drawText(text_obj)
            
            elif hasattr(shape, "image") and shape.image:
                try:
                    img_data = shape.image.blob
                    img_io = io.BytesIO(img_data)
                    c.drawImage(canvas.ImageReader(img_io), left, slide_height - top - height, width, height)
                except Exception as img_err:
                    print(f"Error rendering PPTX shape image: {img_err}")
                    
        c.showPage()
    c.save()

def convert_pdf_to_pptx(pdf_path, pptx_path):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    doc = fitz.open(pdf_path)
    blank_layout = prs.slide_layouts[6]
    
    for page in doc:
        slide = prs.slides.add_slide(blank_layout)
        rect = page.rect
        pdf_w, pdf_h = rect.width, rect.height
        
        blocks = page.get_text("blocks")
        blocks.sort(key=lambda b: (b[1], b[0]))
        
        for block in blocks:
            x0, y0, x1, y1, text, block_no, block_type = block
            block_text = text.strip()
            
            if block_text:
                scale_x = 10.0 / pdf_w
                scale_y = 7.5 / pdf_h
                
                left = Inches(x0 * scale_x)
                top = Inches(y0 * scale_y)
                width = Inches((x1 - x0) * scale_x)
                height = Inches((y1 - y0) * scale_y)
                
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.word_wrap = True
                
                lines = block_text.split('\n')
                for idx, line in enumerate(lines):
                    if idx == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = line
                    p.font.size = Pt(12)
                    p.font.name = 'Helvetica'
        
        try:
            image_list = page.get_images(full=True)
            for img_idx, img in enumerate(image_list):
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                
                rects = page.get_image_rects(xref)
                if rects:
                    img_rect = rects[0]
                    img_x0, img_y0, img_x1, img_y1 = img_rect
                    
                    scale_x = 10.0 / pdf_w
                    scale_y = 7.5 / pdf_h
                    
                    left = Inches(img_x0 * scale_x)
                    top = Inches(img_y0 * scale_y)
                    width = Inches((img_x1 - img_x0) * scale_x)
                    height = Inches((img_y1 - img_y0) * scale_y)
                    
                    img_io = io.BytesIO(image_bytes)
                    slide.shapes.add_picture(img_io, left, top, width, height)
        except Exception as img_err:
            print(f"Error adding slide picture: {img_err}")
            
    prs.save(pptx_path)
    doc.close()

def convert_excel_to_pdf(excel_path, pdf_path):
    wb = openpyxl.load_workbook(excel_path, data_only=True)
    doc = SimpleDocTemplate(pdf_path, pagesize=landscape(letter),
                            rightMargin=36, leftMargin=36,
                            topMargin=36, bottomMargin=36)
    
    styles = getSampleStyleSheet()
    normal_style = ParagraphStyle(
        name='ExcelNormal',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=9,
        leading=11
    )
    title_style = ParagraphStyle(
        name='ExcelTitle',
        parent=styles['Heading1'],
        fontName='Helvetica-Bold',
        fontSize=14,
        leading=18,
        spaceAfter=10
    )
    
    story = []
    
    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        story.append(Paragraph(f"Sheet: {sheet_name}", title_style))
        story.append(Spacer(1, 10))
        
        table_data = []
        max_col = sheet.max_column
        max_row = sheet.max_row
        
        if max_row == 1 and max_col == 1 and sheet.cell(1, 1).value is None:
            story.append(Paragraph("Empty sheet.", normal_style))
            story.append(Spacer(1, 15))
            continue
            
        limit_rows = min(max_row, 200)
        limit_cols = min(max_col, 20)
        
        for r in range(1, limit_rows + 1):
            row_data = []
            for c in range(1, limit_cols + 1):
                val = sheet.cell(row=r, column=c).value
                val_str = str(val) if val is not None else ""
                row_data.append(Paragraph(val_str, normal_style))
            table_data.append(row_data)
            
        if table_data:
            rl_table = RLTable(table_data, colWidths=None)
            rl_table.setStyle(RLTableStyle([
                ('BACKGROUND', (0,0), (-1,0), colors.lightgrey),
                ('ALIGN', (0,0), (-1,-1), 'LEFT'),
                ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
                ('GRID', (0,0), (-1,-1), 0.5, colors.grey),
                ('TOPPADDING', (0,0), (-1,-1), 4),
                ('BOTTOMPADDING', (0,0), (-1,-1), 4),
                ('LEFTPADDING', (0,0), (-1,-1), 4),
                ('RIGHTPADDING', (0,0), (-1,-1), 4),
            ]))
            story.append(rl_table)
            story.append(Spacer(1, 15))
            
        story.append(PageBreak())
        
    if story and isinstance(story[-1], PageBreak):
        story.pop()
        
    doc.build(story)

def convert_pdf_to_excel(pdf_path, excel_path):
    doc = fitz.open(pdf_path)
    wb = openpyxl.Workbook()
    wb.remove(wb.active)
    
    for idx, page in enumerate(doc):
        ws = wb.create_sheet(title=f"Page {idx+1}")
        words = page.get_text("words")
        
        rows = {}
        for w in words:
            y0 = w[1]
            row_key = round(y0 / 6.0) * 6.0
            if row_key not in rows:
                rows[row_key] = []
            rows[row_key].append(w)
            
        sorted_row_keys = sorted(rows.keys())
        excel_row_idx = 1
        
        for r_key in sorted_row_keys:
            row_words = sorted(rows[r_key], key=lambda w: w[0])
            cells = []
            current_cell_text = []
            last_x1 = None
            
            for w in row_words:
                x0, y0, x1, y1, text = w[0], w[1], w[2], w[3], w[4]
                if last_x1 is not None and (x0 - last_x1) > 15:
                    cells.append(" ".join(current_cell_text))
                    current_cell_text = []
                current_cell_text.append(text)
                last_x1 = x1
            if current_cell_text:
                cells.append(" ".join(current_cell_text))
                
            for c_idx, cell_val in enumerate(cells):
                ws.cell(row=excel_row_idx, column=c_idx+1, value=cell_val)
            excel_row_idx += 1
            
    if not wb.sheetnames:
        wb.create_sheet(title="Sheet 1")
        
    wb.save(excel_path)
    doc.close()

def run_ocr_on_pdf(pdf_path, output_pdf_path):
    try:
        import pytesseract
        pytesseract.get_tesseract_version()
        has_tesseract = True
    except Exception:
        has_tesseract = False
        
    doc = fitz.open(pdf_path)
    
    for page in doc:
        pix = page.get_pixmap(dpi=150)
        img_data = pix.tobytes("png")
        img = Image.open(io.BytesIO(img_data))
        
        if has_tesseract:
            try:
                ocr_data = pytesseract.image_to_data(img, output_type=pytesseract.Output.DICT)
                for i in range(len(ocr_data['text'])):
                    text = ocr_data['text'][i].strip()
                    if text:
                        x = ocr_data['left'][i] * 72.0 / 150.0
                        y = ocr_data['top'][i] * 72.0 / 150.0
                        w = ocr_data['width'][i] * 72.0 / 150.0
                        h = ocr_data['height'][i] * 72.0 / 150.0
                        rect = fitz.Rect(x, y, x + w, y + h)
                        page.insert_text(rect.bl, text, fontsize=h, color=(0,0,0), render_mode=3)
            except Exception as ocr_err:
                print(f"OCR execution failure: {ocr_err}")
                text = page.get_text()
                if text:
                    page.insert_text((50, 50), text, fontsize=10, render_mode=3)
        else:
            text = page.get_text()
            if text:
                page.insert_text((50, 50), text, fontsize=10, render_mode=3)
                
    doc.save(output_pdf_path)
    doc.close()

def translate_pdf_document(input_path, output_path, target_lang):
    doc = fitz.open(input_path)
    translator = GoogleTranslator(source='auto', target=target_lang)
    
    pdf_doc = SimpleDocTemplate(output_path, pagesize=letter)
    styles = getSampleStyleSheet()
    normal_style = ParagraphStyle(
        name='TransNormal',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=11,
        leading=14,
        spaceAfter=10
    )
    
    story = []
    for page in doc:
        text = page.get_text("text").strip()
        if text:
            paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()]
            for p in paragraphs:
                if len(p) > 4000:
                    chunks = [p[i:i+4000] for i in range(0, len(p), 4000)]
                else:
                    chunks = [p]
                    
                for chunk in chunks:
                    try:
                        translated = translator.translate(chunk)
                        safe_text = translated.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                        story.append(Paragraph(safe_text, normal_style))
                    except Exception:
                        safe_text = chunk.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                        story.append(Paragraph(safe_text, normal_style))
            story.append(Spacer(1, 15))
            
    doc.close()
    if not story:
        story.append(Paragraph("Empty or non-extractable PDF document.", normal_style))
        
    pdf_doc.build(story)

def summarize_pdf_document(input_path, output_path):
    import re
    from collections import Counter
    
    doc = fitz.open(input_path)
    text = ""
    for page in doc:
        text += page.get_text() + "\n"
    doc.close()
    
    text = text.strip()
    if not text:
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write("No extractable text found to summarize.")
        return
        
    sentences = re.split(r'(?<!\w\.\w.)(?<![A-Z][a-z]\.)(?<=\.|\?)\s', text)
    sentences = [s.strip() for s in sentences if len(s.strip()) > 15]
    
    if len(sentences) <= 5:
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(text)
        return
        
    words = re.findall(r'\b\w+\b', text.lower())
    stopwords = set([
        'the', 'a', 'and', 'is', 'of', 'in', 'to', 'for', 'it', 'on', 'with', 'as', 
        'at', 'by', 'an', 'be', 'this', 'are', 'from', 'that', 'or', 'but', 'not', 
        'he', 'she', 'they', 'we', 'you', 'i', 'has', 'have', 'had', 'was', 'were',
        'been', 'will', 'would', 'should', 'can', 'could', 'may', 'might', 'must',
        'about', 'into', 'than', 'then', 'them', 'their', 'there', 'who', 'which',
        'what', 'when', 'where', 'how', 'why', 'all', 'any', 'both', 'each', 'few',
        'more', 'most', 'other', 'some', 'such', 'no', 'nor', 'too', 'very'
    ])
    filtered_words = [w for w in words if w not in stopwords and len(w) > 2]
    
    word_freq = Counter(filtered_words)
    max_freq = max(word_freq.values()) if word_freq else 1
    
    for word in word_freq:
        word_freq[word] = word_freq[word] / max_freq
        
    sentence_scores = {}
    for sent in sentences:
        score = 0
        sent_words = re.findall(r'\b\w+\b', sent.lower())
        for w in sent_words:
            if w in word_freq:
                score += word_freq[w]
        sentence_scores[sent] = score
        
    num_sentences = max(5, min(8, len(sentences) // 5))
    top_sentences = sorted(sentence_scores.keys(), key=lambda s: sentence_scores[s], reverse=True)[:num_sentences]
    top_sentences.sort(key=lambda s: sentences.index(s))
    
    summary = "\n\n".join([f"- {s}" for s in top_sentences])
    
    original_word_count = len(text.split())
    summary_word_count = len(" ".join(top_sentences).split())
    savings = round((1 - summary_word_count / original_word_count) * 100)
    
    summary_header = f"=== AI DOCUMENT SUMMARY ===\n" \
                     f"Original length: {original_word_count} words\n" \
                     f"Summary length: {summary_word_count} words (saved {savings}% of reading time)\n" \
                     f"===========================\n\n"
                     
    with open(output_path, 'w', encoding='utf-8') as f:
        f.write(summary_header + summary)
